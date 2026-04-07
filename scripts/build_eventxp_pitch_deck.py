#!/usr/bin/env python3
"""Generate EventXP client pitch deck (pptx). Run: python3 scripts/build_eventxp_pitch_deck.py"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


def add_title_slide(prs, title: str, subtitle: str = ""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    left, top, width, height = Inches(0.5), Inches(1.8), Inches(9), Inches(1.2)
    tx = slide.shapes.add_textbox(left, top, width, height)
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    if subtitle:
        st = tx.text_frame.add_paragraph()
        st.text = subtitle
        st.font.size = Pt(18)
        st.font.color.rgb = RGBColor(71, 85, 105)
        st.space_before = Pt(12)


def add_bullet_slide(prs, title: str, bullets: list[str], footer: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Title
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.45), Inches(9), Inches(0.9))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(15, 23, 42)

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(9), Inches(4.8))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(8)
        p.level = 0

    if footer:
        fb = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.55))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(11)
        fp.font.color.rgb = RGBColor(100, 116, 139)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "EventXP",
        "AI Community Growth Engine · Client Pitch\nInnovateXP Limited · https://innovatexp.co",
    )

    add_bullet_slide(
        prs,
        "Executive summary",
        [
            "將「活動」變成可追蹤、可複製、可增長的社群資產。",
            "定位：社群增長 + 轉化（Leads / Retention / ROI），唔係只做報名。",
            "交付：平台整合 + AI 洞察 + 落地支援（現場 + 策略諮詢）。",
        ],
    )

    add_bullet_slide(
        prs,
        "痛點（Problems）",
        [
            "活動後數據散、匯出慢、缺席/嘉賓狀態唔清晰。",
            "跟進無優先序：唔知邊啲人最值開 Conversation、邊啲會回流。",
            "ROI 難向董事局/贊助商證明：得出席，冇增長故事。",
        ],
    )

    add_bullet_slide(
        prs,
        "產品詳細介紹（Modules）",
        [
            "Luma / 平台整合：名單同步、流程銜接（少改現有流程）。",
            "Smart Conversion List：高潛力 leads 自動整理，方便跟進。",
            "Retention Insights：回流與出席型態分析（數據累積後更準）。",
            "Reporting & Export：即時出席、會員缺席、嘉賓；CSV 與營運對賬。",
            "Note：進階「關單率」建議以可選 CRM 整合表述，避免 overpromise。",
        ],
    )

    add_bullet_slide(
        prs,
        "定位（Positioning）",
        [
            "專為專業社群 / 商務網絡：要增長、要轉化、要留存。",
            "唔係純票務平台：補足「執行」與「數據變現」之間嘅空白。",
            "中型社群高性價比方案：比單純工具多一層增長與落地。",
        ],
    )

    add_bullet_slide(
        prs,
        "成果與 KPI（What we measure）",
        [
            "出席率、缺席、嘉賓到場與登記。",
            "回流/出席型態、邀請來源（視數據可用性）。",
            "轉化名單數量與跟進速度（Sales motion 可選）。",
        ],
    )

    add_bullet_slide(
        prs,
        "服務與上線（Delivery）",
        [
            "Go-live：整合設定、測試、培訓與文件。",
            "現場支援：按套餐場次；確保活動日無斷線。",
            "策略諮詢：Starter/ Growth / Enterprise 分層。",
        ],
    )

    add_bullet_slide(
        prs,
        "建議定價（Suggest Pricing）",
        [
            "一次性 Setup：約 HKD 14,000（整合 + 初始設定 + 文件/培訓）。",
            "Starter：HKD 1,480/月 — 1 社群；基本報告與匯出。",
            "Growth：HKD 1,980/月 — 2 社群；AI Insights + Leads；年內 2 次諮詢 + 3 場現場。",
            "Enterprise：HKD 2,480/月 起 — 4+ 社群；月諮詢 + 4 場現場；可選深度整合。",
            "預付：1 年 -10%；2 年 -20%（按現金流協定）。",
        ],
    )

    add_bullet_slide(
        prs,
        "Commission / Partner package",
        [
            "Referral / Reseller：建議 25%（月費首 12 個月，或長期 15% 二選一）。",
            "一次性 Setup：可設 10%–15%（視伙伴參與度）。",
            "條款：付款成功計佣；退款回撈；分工：引薦 vs closing vs 交付可寫明。",
        ],
    )

    add_bullet_slide(
        prs,
        "SWOT（高層摘要）",
        [
            "S：AI+營運一體化；整合現有平台；落地支援強。",
            "W：第三方 API 依賴；AI 需數據累積期。",
            "O：專業社群數據化；中型市場缺「增長工具」；伙伴網絡擴張。",
            "T：競品插件；活動預算波動；PDPO/私隱合規成本上升。",
        ],
    )

    add_bullet_slide(
        prs,
        "Why EventXP？",
        [
            "由「搞掂活動」升級到「搞掂增長」：邊個值得跟進、點樣提升回流。",
            "唔使大改流程：係 upgrade，唔係 replace。",
            "現場 + 策略：確保用得成，唔係買完唔用。",
            "共贏：伙伴分潤，一齊放大影響力。",
        ],
        footer="Contact: https://innovatexp.co",
    )

    add_bullet_slide(
        prs,
        "Next steps（CTA）",
        [
            "確認範圍：社群數量、平台（如 Luma）、場次與合规要求。",
            "簽署報價 / SOW → 開通整合 → UAT → Go-live。",
            "聯絡我們：innovatexp.co",
        ],
    )

    out = Path(__file__).resolve().parents[1] / "docs" / "EventXP_Client_Pitch.pptx"
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(out)


if __name__ == "__main__":
    main()
