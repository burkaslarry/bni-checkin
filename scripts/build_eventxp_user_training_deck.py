#!/usr/bin/env python3
"""Generate EventXP end-user training deck (Traditional Chinese). Run from repo root:

    pip install python-pptx   # if needed
    python3 scripts/build_eventxp_user_training_deck.py

Output: docs/EventXP_User_Training_ZH.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "docs" / "EventXP_User_Training_ZH.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(9), Inches(1.4))
    p = tx.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor(15, 23, 42)
    if subtitle:
        st = tx.text_frame.add_paragraph()
        st.text = subtitle
        st.font.size = Pt(17)
        st.font.color.rgb = RGBColor(71, 85, 105)
        st.space_before = Pt(10)


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    footer: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbox = slide.shapes.add_textbox(Inches(0.55), Inches(0.45), Inches(9), Inches(0.85))
    tp = tbox.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(26)
    tp.font.bold = True
    tp.font.color.rgb = RGBColor(15, 23, 42)

    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.25), Inches(9), Inches(5.1))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = line
        para.font.size = Pt(15)
        para.font.color.rgb = RGBColor(30, 41, 59)
        para.space_after = Pt(7)

    if footer:
        fb = slide.shapes.add_textbox(Inches(0.55), Inches(6.55), Inches(9), Inches(0.6))
        fp = fb.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(11)
        fp.font.color.rgb = RGBColor(100, 116, 139)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    urls = (
        "簽到首頁：https://bni-anchor-checkin.vercel.app\n"
        "管理後台：https://bni-anchor-checkin.vercel.app/admin"
    )

    add_title_slide(
        prs,
        "EventXP 使用教學",
        "BNI Anchor 活動簽到 · 現場＋管理員操作\nInnovateXP Limited · innovatexp.co",
    )

    add_bullet_slide(
        prs,
        "網址同三大角色",
        [
            urls,
            "會員：掃描活動 QR，或由列表揀姓名簽到。",
            "嘉賓（預先登記）：喺簽到表單用搜尋揀自己，再確認簽到。",
            "管理員：後台建立／切換「當前活動」、產 QR、睇報告、匯出 CSV。",
        ],
    )

    add_bullet_slide(
        prs,
        "會員簽到（現場）",
        [
            "1. 打開簽到首頁 → 進入會員簽到。",
            "2. 用相機掃描**當次活動** QR（或由主辦提供連結進入）。",
            "3. 若掃描有問題：可用列表搜尋姓名／專業，揀中後按確認簽到。",
            "4. 已成功簽到會提示；重複簽到會被拒絕。",
        ],
        footer="提示：簽到綁定「當前活動」，現場請確認主辦已設好活動。",
    )

    add_bullet_slide(
        prs,
        "嘉賓簽到（預先登記名單）",
        [
            "1. 簽到頁揀「嘉賓」→ 用搜尋框搵自己姓名或專業。",
            "2. 揀中後按確認簽到。",
            "3. 名冊可按右側圖示重新載入，與後台同步。",
        ],
    )

    add_bullet_slide(
        prs,
        "公開嘉賓登記（無需登入）",
        [
            "連結形如：…/public/guest?eventID=（活動編號）",
            "管理員可喺後台「公開嘉賓連結」頁複製連結分享。",
            "填寫必填欄位＋驗證題答案正確後，才能提交（防機械人）。",
            "提交後資料入庫，可喺「嘉賓管理」用活動日期篩選查看。",
        ],
    )

    add_bullet_slide(
        prs,
        "管理後台 · 活動同 QR",
        [
            "進入 /admin →「新增活動和二維碼」。",
            "填活動名稱、日期、時段、準時截止時間等 → 建立活動。",
            "喺「活動管理」將正進行緊嘅活動設為「當前活動」（現場簽到會跟呢個）。",
            "產生／下載 QR，投影或印刷俾會員掃描。",
        ],
    )

    add_bullet_slide(
        prs,
        "報告、記錄、匯出 CSV",
        [
            "「即時簽到狀態」頁（/report）：出席／缺席、篩選、自動更新。",
            "可切換「簽到記錄 CSV」分頁：搜尋、篩選、單筆刪除、匯出。",
            "後台「匯出」區：可從伺服器揀檔名下載（已對齊活動嘅出席狀態）。",
        ],
    )

    add_bullet_slide(
        prs,
        "常見情況",
        [
            "掃描失敗 → 改用手動揀會員／確認光線用緊正確活動 QR。",
            "顯示已簽到 → 每人每活動一般只可簽一次；有需要搵管理員。",
            "匯出亂碼 → 檔案為 UTF-8 BOM；用 Excel 開啟通常正常。",
            "會員／嘉賓名單來自後端資料庫；更新請聯絡管理員。",
        ],
    )

    add_bullet_slide(
        prs,
        "需要幫手？",
        [
            "InnovateXP Limited",
            "https://innovatexp.co",
            "",
            "本簡報可由 repo 內 `python3 scripts/build_eventxp_user_training_deck.py` 重新產生。",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
